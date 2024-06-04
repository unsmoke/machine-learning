from pptx import Presentation
from pptx.util import Inches

# Membuat objek presentasi baru
prs = Presentation()

# Fungsi untuk menambah slide dengan judul dan konten
def add_slide(prs, title, content, layout=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

# Slide 1: Judul Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "Pitch Deck - Unsmoke"
subtitle_1.text = "Perlombaan Startup Internasional"

# Slide 2: Masalah yang Dihadapi
add_slide(prs, "The Problem", 
          "• Tingginya angka perokok di Indonesia\n"
          "• Rendahnya tingkat keberhasilan berhenti merokok\n"
          "• Dampak kesehatan yang serius dari merokok")

# Slide 3: Solusi yang Ditawarkan
add_slide(prs, "Our Solution", 
          "• Unsmoke: Aplikasi yang menggunakan gamifikasi untuk membantu berhenti merokok\n"
          "• Penggunaan AI untuk personalisasi rencana berhenti merokok\n"
          "• Fitur-fitur inovatif yang memotivasi pengguna")

# Slide 4: Produk dan Fitur Utama
add_slide(prs, "Product and Key Features", 
          "• Pelacakan kemajuan dan manfaat\n"
          "• Animasi avatar paru-paru\n"
          "• Tantangan interaktif dan leaderboard\n"
          "• Toko makeover paru-paru\n"
          "• Rencana berhenti yang dipersonalisasi dengan AI\n"
          "• Pelacakan kesehatan harian")

# Slide 5: Analisis Pasar
add_slide(prs, "Market Analysis", 
          "• Ukuran pasar potensial (TAM, SAM, SOM)\n"
          "• Target pasar utama: Generasi Z di Indonesia\n"
          "• Data statistik terkait perokok di Indonesia")

# Slide 6: Model Bisnis
add_slide(prs, "Business Model", 
          "• Langganan premium\n"
          "• Produk dan layanan tambahan\n"
          "• Kemitraan dengan perusahaan asuransi kesehatan")

# Slide 7: Strategi Go-to-Market
add_slide(prs, "Go-to-Market Strategy", 
          "• Kampanye pemasaran digital\n"
          "• Kolaborasi dengan penyedia layanan telekomunikasi\n"
          "• Kampanye edukasi di wilayah dengan penetrasi internet rendah")

# Slide 8: Analisis Kompetitor
add_slide(prs, "Competitor Analysis", 
          "• Perbandingan fitur dengan kompetitor utama\n"
          "• Keunggulan kompetitif Unsmoke")

# Slide 9: Proyeksi Keuangan
add_slide(prs, "Financial Projections", 
          "• Proyeksi pendapatan dari langganan premium\n"
          "• Estimasi pengeluaran dan keuntungan")

# Slide 10: Tim Inti
add_slide(prs, "Core Team", 
          "• Anggota tim dan peran masing-masing\n"
          "• Pengalaman dan keahlian yang relevan")

# Slide 11: Rencana Pengembangan Produk
add_slide(prs, "Product Development Plan", 
          "• Timeline pengembangan produk\n"
          "• Fitur-fitur yang akan datang")

# Slide 12: Kesimpulan dan Ajakan Bertindak
add_slide(prs, "Conclusion and Call to Action", 
          "• Ringkasan nilai tambah Unsmoke\n"
          "• Ajakan untuk investasi atau kolaborasi")

# Menyimpan presentasi
prs.save('Unsmoke_Pitch_Deck_Final.pptx')
